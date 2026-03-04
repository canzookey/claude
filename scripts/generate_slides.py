#!/usr/bin/env python3
"""
株式会社ビジコン - スライド生成ツール
テンプレートを使って新しいスライドや提案資料を生成します。
"""

import argparse
import copy
import os
import sys
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor

TEMPLATE_PATH = os.path.join(os.path.dirname(__file__), "..", "templates", "【ご紹介資料】株式会社ビジコン.pptx")
OUTPUT_DIR = os.path.join(os.path.dirname(__file__), "..", "output")


def load_template():
    """テンプレートを読み込む"""
    if not os.path.exists(TEMPLATE_PATH):
        print(f"エラー: テンプレートが見つかりません: {TEMPLATE_PATH}")
        sys.exit(1)
    return Presentation(TEMPLATE_PATH)


def find_case_study_slide(prs, index=7):
    """事例スライドを取得（デフォルト: スライド8 = index 7）"""
    if index < len(prs.slides):
        return prs.slides[index]
    return None


def update_text_in_shape(shape, old_text, new_text):
    """シェイプ内のテキストを置換"""
    if not shape.has_text_frame:
        return False
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if old_text in run.text:
                run.text = run.text.replace(old_text, new_text)
                return True
    return False


def set_shape_text(shape, text):
    """シェイプのテキストを設定（書式を保持）"""
    if not shape.has_text_frame:
        return
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            run.text = ""
    if shape.text_frame.paragraphs and shape.text_frame.paragraphs[0].runs:
        shape.text_frame.paragraphs[0].runs[0].text = text


def duplicate_slide(prs, slide_index):
    """スライドを複製する"""
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT
    import copy
    from lxml import etree

    template_slide = prs.slides[slide_index]
    slide_layout = template_slide.slide_layout

    new_slide = prs.slides.add_slide(slide_layout)

    # 既存の要素をクリア
    for shape in list(new_slide.shapes):
        sp = shape._element
        sp.getparent().remove(sp)

    # テンプレートスライドの要素をコピー
    for shape in template_slide.shapes:
        el = copy.deepcopy(shape._element)
        new_slide.shapes._spTree.append(el)

    return new_slide


def add_case_study(args):
    """事例スライドを追加"""
    prs = load_template()

    # スライド8（index=7）を事例テンプレートとして使用
    ref_slide_index = 7
    new_slide = duplicate_slide(prs, ref_slide_index)

    # テキスト置換のマッピング
    replacements = {
        "事例：DX推進室伴走支援（老舗カステラ製造会社様）": args.title,
        "製造業": args.industry,
    }

    for shape in new_slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            full_text = para.text.strip()

            # タイトル
            if "事例：" in full_text or "事例:" in full_text:
                for run in para.runs:
                    if "事例" in run.text:
                        run.text = args.title
                        break

            # サブタイトル（最初の説明文）
            if "資材管理・動画マニュアル" in full_text:
                for run in para.runs:
                    run.text = ""
                para.runs[0].text = args.subtitle

            # 業種タグ
            if full_text == "製造業":
                for run in para.runs:
                    run.text = args.industry

            # 課題
            if "自社プロダクトのリブランディング" in full_text or "業務フロー作成" in full_text:
                for run in para.runs:
                    run.text = ""
                para.runs[0].text = args.challenge

            # 支援内容
            if "デジタルロードマップ" in full_text or "ITツールのベンダー選定" in full_text:
                for run in para.runs:
                    run.text = ""
                para.runs[0].text = args.support

            # 成果
            if "ITツール導入を導入する" in full_text:
                for run in para.runs:
                    run.text = ""
                para.runs[0].text = args.result

    # 保存
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    output_path = os.path.join(OUTPUT_DIR, f"case_study_{timestamp}.pptx")
    prs.save(output_path)
    print(f"事例スライドを追加しました: {output_path}")
    return output_path


def new_proposal(args):
    """新しい提案資料を作成"""
    prs = load_template()

    # 表紙を更新
    cover_slide = prs.slides[0]
    for shape in cover_slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                if "ご紹介資料" in para.text:
                    for run in para.runs:
                        if "ご紹介資料" in run.text:
                            run.text = args.theme
                            break

    # 保存
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    safe_client = args.client.replace("/", "_").replace("\\", "_")
    output_path = os.path.join(OUTPUT_DIR, f"proposal_{safe_client}_{timestamp}.pptx")
    prs.save(output_path)
    print(f"提案資料を作成しました: {output_path}")
    return output_path


def create_case_only(args):
    """事例スライドのみの資料を新規作成"""
    prs = Presentation()
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)

    # テンプレートから事例スライドのデザインを参照
    template = load_template()

    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # テキストボックスを追加
    from pptx.util import Pt, Inches

    # タイトル
    txBox = slide.shapes.add_textbox(Emu(495088), Emu(366960), Emu(9739883), Emu(369332))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = args.title
    p.font.size = Pt(18)
    p.font.bold = True

    # サブタイトル
    txBox = slide.shapes.add_textbox(Emu(376749), Emu(1279771), Emu(5903987), Emu(800219))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = args.subtitle
    p.font.size = Pt(12)

    # 業種タグ
    txBox = slide.shapes.add_textbox(Emu(10234972), Emu(296062), Emu(1461939), Emu(495135))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = args.industry
    p.font.size = Pt(11)

    # 課題ラベル
    txBox = slide.shapes.add_textbox(Emu(346894), Emu(2468956), Emu(414596), Emu(964762))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "課題"
    p.font.size = Pt(14)
    p.font.bold = True

    # 課題内容
    txBox = slide.shapes.add_textbox(Emu(936762), Emu(2474479), Emu(4604134), Emu(1120307))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = args.challenge
    p.font.size = Pt(10)

    # 支援内容ラベル
    txBox = slide.shapes.add_textbox(Emu(346893), Emu(4173710), Emu(414596), Emu(825130))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "支援内容"
    p.font.size = Pt(14)
    p.font.bold = True

    # 支援内容
    txBox = slide.shapes.add_textbox(Emu(936761), Emu(3642224), Emu(4604134), Emu(1400383))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = args.support
    p.font.size = Pt(10)

    # 成果ラベル
    txBox = slide.shapes.add_textbox(Emu(346893), Emu(5764718), Emu(414596), Emu(964762))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "成果"
    p.font.size = Pt(14)
    p.font.bold = True

    # 成果内容
    txBox = slide.shapes.add_textbox(Emu(936761), Emu(5770232), Emu(4604134), Emu(840230))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = args.result
    p.font.size = Pt(10)

    # 保存
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    output_path = os.path.join(OUTPUT_DIR, f"case_{timestamp}.pptx")
    prs.save(output_path)
    print(f"事例スライドを作成しました: {output_path}")
    return output_path


def list_slides(args):
    """テンプレートのスライド一覧を表示"""
    prs = load_template()
    print(f"テンプレート: {os.path.basename(TEMPLATE_PATH)}")
    print(f"スライド数: {len(prs.slides)}")
    print()
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text and len(text) > 2:
                        texts.append(text)
        title = texts[0] if texts else "(テキストなし)"
        print(f"  スライド {i+1}: {title[:60]}")


def main():
    parser = argparse.ArgumentParser(description="ビジコン スライド生成ツール")
    subparsers = parser.add_subparsers(dest="command", help="コマンド")

    # add-case: 事例スライド追加
    case_parser = subparsers.add_parser("add-case", help="テンプレートに事例スライドを追加")
    case_parser.add_argument("--title", required=True, help="事例タイトル（例：事例：〇〇（△△会社様））")
    case_parser.add_argument("--subtitle", required=True, help="成果の要約（1-2行）")
    case_parser.add_argument("--industry", required=True, help="業種タグ（例：製造業、建設業）")
    case_parser.add_argument("--challenge", required=True, help="課題の説明")
    case_parser.add_argument("--support", required=True, help="支援内容の説明")
    case_parser.add_argument("--result", required=True, help="成果の説明")

    # new-proposal: 提案資料作成
    prop_parser = subparsers.add_parser("new-proposal", help="新しい提案資料を作成")
    prop_parser.add_argument("--client", required=True, help="クライアント名")
    prop_parser.add_argument("--theme", required=True, help="提案テーマ")

    # create-case: 事例スライドのみ新規作成
    create_parser = subparsers.add_parser("create-case", help="事例スライドのみの資料を新規作成")
    create_parser.add_argument("--title", required=True, help="事例タイトル")
    create_parser.add_argument("--subtitle", required=True, help="成果の要約")
    create_parser.add_argument("--industry", required=True, help="業種タグ")
    create_parser.add_argument("--challenge", required=True, help="課題の説明")
    create_parser.add_argument("--support", required=True, help="支援内容の説明")
    create_parser.add_argument("--result", required=True, help="成果の説明")

    # list: スライド一覧
    subparsers.add_parser("list", help="テンプレートのスライド一覧を表示")

    args = parser.parse_args()

    if args.command == "add-case":
        add_case_study(args)
    elif args.command == "new-proposal":
        new_proposal(args)
    elif args.command == "create-case":
        create_case_only(args)
    elif args.command == "list":
        list_slides(args)
    else:
        parser.print_help()


if __name__ == "__main__":
    main()
