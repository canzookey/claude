#!/usr/bin/env python3
"""PPTXスライドをPNG画像としてレンダリングするスクリプト"""

from pptx import Presentation
from pptx.util import Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image, ImageDraw, ImageFont
import os

SCALE = 2
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5
DPI = 96
SLIDE_W = int(SLIDE_W_IN * DPI * SCALE)
SLIDE_H = int(SLIDE_H_IN * DPI * SCALE)

FONT_REGULAR = "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc"
FONT_BOLD = "/usr/share/fonts/opentype/noto/NotoSansCJK-Bold.ttc"

_font_cache = {}

def get_font(size_pt, bold=False):
    key = (size_pt, bold)
    if key in _font_cache:
        return _font_cache[key]
    px = int(size_pt * SCALE)
    path = FONT_BOLD if bold else FONT_REGULAR
    try:
        f = ImageFont.truetype(path, px, index=0)
    except:
        f = ImageFont.load_default()
    _font_cache[key] = f
    return f

def emu_to_px(emu_val):
    return int(emu_val / 914400 * DPI * SCALE)

def parse_rgb(color_obj):
    try:
        if color_obj and color_obj.rgb:
            h = str(color_obj.rgb)
            return (int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))
    except:
        pass
    return None

def get_fill(shape):
    try:
        f = shape.fill
        if f.type is not None:
            return parse_rgb(f.fore_color)
    except:
        pass
    return None

def get_cell_fill(cell):
    try:
        f = cell.fill
        if f.type is not None:
            return parse_rgb(f.fore_color)
    except:
        pass
    return None

def get_bg(slide):
    try:
        f = slide.background.fill
        if f.type is not None:
            return parse_rgb(f.fore_color)
    except:
        pass
    return (255,255,255)

def wrap_text(text, font, max_w):
    if not text:
        return [""]
    lines = []
    for raw_line in text.split('\n'):
        if not raw_line:
            lines.append("")
            continue
        current = ""
        for ch in raw_line:
            test = current + ch
            bbox = font.getbbox(test)
            if bbox[2] - bbox[0] > max_w and current:
                lines.append(current)
                current = ch
            else:
                current = test
        lines.append(current)
    return lines

def draw_text_block(draw, text, x, y, w, h, font_size, color, bold, align):
    if not text or not text.strip():
        return 0
    font = get_font(font_size, bold)
    lines = wrap_text(text, font, w)
    line_h = int(font_size * SCALE) + 6 * SCALE
    cy = y
    for line in lines:
        if cy + line_h > y + h + 2*SCALE:
            break
        if not line:
            cy += line_h // 2
            continue
        bbox = font.getbbox(line)
        tw = bbox[2] - bbox[0]
        if align == PP_ALIGN.CENTER:
            tx = x + (w - tw) // 2
        elif align == PP_ALIGN.RIGHT:
            tx = x + w - tw
        else:
            tx = x
        draw.text((tx, cy), line, fill=color, font=font)
        cy += line_h
    return cy - y

def render_textframe(draw, tf, x, y, w, h):
    cy = y
    for para in tf.paragraphs:
        text = para.text
        # Determine formatting
        font_size = 14
        color = (0x1B, 0x1B, 0x2F)
        bold = False
        align = para.alignment if para.alignment is not None else PP_ALIGN.LEFT

        # Paragraph-level font
        if para.font and para.font.size:
            font_size = para.font.size.pt
        pfc = parse_rgb(para.font.color) if para.font else None
        if pfc:
            color = pfc
        if para.font and para.font.bold:
            bold = True

        # Run-level overrides
        if para.runs:
            r = para.runs[0]
            if r.font.size:
                font_size = r.font.size.pt
            rfc = parse_rgb(r.font.color)
            if rfc:
                color = rfc
            if r.font.bold:
                bold = True

        if not text:
            cy += int(font_size * SCALE * 0.4)
            continue

        used = draw_text_block(draw, text, x, cy, w, y + h - cy, font_size, color, bold, align)
        cy += max(used, int(font_size * SCALE) + 6*SCALE)

        # space_after
        try:
            sa = para.space_after
            if sa:
                cy += emu_to_px(sa)
        except:
            pass

def render_table(draw, shape):
    left = emu_to_px(shape.left)
    top = emu_to_px(shape.top)
    width = emu_to_px(shape.width)
    height = emu_to_px(shape.height)

    table = shape.table
    n_rows = len(table.rows)
    n_cols = len(table.columns)

    col_widths = [emu_to_px(table.columns[c].width) for c in range(n_cols)]
    row_h = height // n_rows

    for r in range(n_rows):
        for c in range(n_cols):
            cx = left + sum(col_widths[:c])
            cy = top + r * row_h
            cw = col_widths[c]
            ch = row_h

            cell = table.cell(r, c)

            # Cell fill
            fc = get_cell_fill(cell)
            if fc:
                draw.rectangle([cx, cy, cx+cw, cy+ch], fill=fc)
            elif r == 0:
                draw.rectangle([cx, cy, cx+cw, cy+ch], fill=(0x2D,0x5B,0xE3))
            elif r % 2 == 0:
                draw.rectangle([cx, cy, cx+cw, cy+ch], fill=(0xF0,0xF2,0xF5))
            else:
                draw.rectangle([cx, cy, cx+cw, cy+ch], fill=(255,255,255))

            # Border
            draw.rectangle([cx, cy, cx+cw, cy+ch], outline=(0xDD,0xDD,0xDD))

            # Cell text
            text = cell.text
            if not text:
                continue

            font_size = 11
            color = (0x1B, 0x1B, 0x2F)
            bold = False
            align = PP_ALIGN.CENTER

            for para in cell.text_frame.paragraphs:
                if para.font and para.font.size:
                    font_size = para.font.size.pt
                pfc = parse_rgb(para.font.color) if para.font else None
                if pfc:
                    color = pfc
                if para.font and para.font.bold:
                    bold = True
                if para.alignment is not None:
                    align = para.alignment
                if para.runs:
                    run = para.runs[0]
                    if run.font.size:
                        font_size = run.font.size.pt
                    rfc = parse_rgb(run.font.color)
                    if rfc:
                        color = rfc
                    if run.font.bold:
                        bold = True

            pad = 6 * SCALE
            draw_text_block(draw, text, cx+pad, cy+pad, cw-2*pad, ch-2*pad,
                          font_size, color, bold, align)


def render_slide(prs, idx, out_path):
    slide = prs.slides[idx]
    bg = get_bg(slide)
    img = Image.new('RGB', (SLIDE_W, SLIDE_H), bg)
    draw = ImageDraw.Draw(img)

    for shape in slide.shapes:
        try:
            sx = emu_to_px(shape.left)
            sy = emu_to_px(shape.top)
            sw = emu_to_px(shape.width)
            sh = emu_to_px(shape.height)

            if shape.has_table:
                render_table(draw, shape)
                continue

            # Draw shape fill (rectangle)
            fc = get_fill(shape)
            if fc:
                draw.rectangle([sx, sy, sx+sw, sy+sh], fill=fc)

            # Draw text
            if shape.has_text_frame:
                pad = 4 * SCALE
                render_textframe(draw, shape.text_frame,
                               sx + pad, sy + pad, sw - 2*pad, sh - 2*pad)
        except Exception as e:
            pass

    # Downscale
    img = img.resize((SLIDE_W // SCALE, SLIDE_H // SCALE), Image.LANCZOS)
    img.save(out_path, "PNG")
    return out_path


prs = Presentation("/home/user/claude/sales-material-analysis.pptx")
out_dir = "/home/user/claude/slide_preview"
os.makedirs(out_dir, exist_ok=True)

for i in range(len(prs.slides)):
    p = render_slide(prs, i, f"{out_dir}/slide_{i+1:02d}.png")
    print(f"Rendered: {p}")
