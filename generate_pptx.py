#!/usr/bin/env python3
"""SaaS事業計画書 分析レポートをPowerPointに出力するスクリプト"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Color Palette ──
COLOR_DARK = RGBColor(0x1B, 0x1B, 0x2F)
COLOR_PRIMARY = RGBColor(0x2D, 0x5B, 0xE3)
COLOR_ACCENT = RGBColor(0x00, 0xB4, 0xD8)
COLOR_RED = RGBColor(0xE0, 0x3E, 0x3E)
COLOR_GREEN = RGBColor(0x2E, 0x9E, 0x5A)
COLOR_ORANGE = RGBColor(0xE8, 0x8D, 0x2A)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_LIGHT_GRAY = RGBColor(0xF0, 0xF2, 0xF5)
COLOR_GRAY = RGBColor(0x6B, 0x72, 0x80)
COLOR_DARK_GRAY = RGBColor(0x37, 0x3D, 0x4E)


def add_bg(slide, color=COLOR_WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=14,
                color=COLOR_DARK, bold=False, alignment=PP_ALIGN.LEFT, font_name="Yu Gothic"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multi_text(slide, left, top, width, height, lines, font_name="Yu Gothic"):
    """lines: list of (text, font_size, color, bold, alignment)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, font_size, color, bold, alignment) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = alignment
        p.space_after = Pt(4)
    return txBox


def make_table(slide, left, top, width, height, rows, col_widths=None):
    """rows: list of list of str. First row = header."""
    n_rows = len(rows)
    n_cols = len(rows[0])
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for r_idx, row in enumerate(rows):
        for c_idx, cell_text in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = cell_text
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(11)
                paragraph.font.name = "Yu Gothic"
                if r_idx == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = COLOR_WHITE
                    paragraph.alignment = PP_ALIGN.CENTER
                else:
                    paragraph.font.color.rgb = COLOR_DARK
                    if c_idx == 0:
                        paragraph.alignment = PP_ALIGN.LEFT
                    else:
                        paragraph.alignment = PP_ALIGN.CENTER

            if r_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_PRIMARY
            elif r_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_LIGHT_GRAY
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_WHITE

    return table_shape


def add_card(slide, left, top, width, height, title, value, sub="", value_color=COLOR_PRIMARY):
    shape = add_shape_rect(slide, left, top, width, height, COLOR_WHITE)
    shape.shadow.inherit = False

    add_textbox(slide, left + Inches(0.2), top + Inches(0.15), width - Inches(0.4), Inches(0.3),
                title, font_size=10, color=COLOR_GRAY, bold=False)
    add_textbox(slide, left + Inches(0.2), top + Inches(0.45), width - Inches(0.4), Inches(0.5),
                value, font_size=24, color=value_color, bold=True)
    if sub:
        add_textbox(slide, left + Inches(0.2), top + Inches(0.95), width - Inches(0.4), Inches(0.3),
                    sub, font_size=9, color=COLOR_GRAY, bold=False)


def slide_header(slide, title, subtitle=""):
    add_shape_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), COLOR_PRIMARY)
    add_textbox(slide, Inches(0.6), Inches(0.3), Inches(10), Inches(0.6),
                title, font_size=26, color=COLOR_DARK, bold=True)
    if subtitle:
        add_textbox(slide, Inches(0.6), Inches(0.8), Inches(10), Inches(0.4),
                    subtitle, font_size=13, color=COLOR_GRAY)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# Slide 1: Title
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_bg(slide, COLOR_DARK)
add_shape_rect(slide, Inches(0), Inches(3.2), Inches(13.333), Inches(0.06), COLOR_ACCENT)

add_textbox(slide, Inches(1), Inches(1.5), Inches(11), Inches(1),
            "SaaS事業計画書", font_size=44, color=COLOR_WHITE, bold=True)
add_textbox(slide, Inches(1), Inches(2.3), Inches(11), Inches(0.6),
            "営業資料 分析レポート", font_size=24, color=COLOR_ACCENT, bold=False)

add_multi_text(slide, Inches(1), Inches(3.8), Inches(11), Inches(2.5), [
    ("中小企業向けクラウドベース業務効率化SaaSプラットフォーム", 16, COLOR_WHITE, False, PP_ALIGN.LEFT),
    ("", 10, COLOR_WHITE, False, PP_ALIGN.LEFT),
    ("分析日: 2026年2月25日", 13, COLOR_GRAY, False, PP_ALIGN.LEFT),
    ("対象資料: saas-business-plan.md", 13, COLOR_GRAY, False, PP_ALIGN.LEFT),
])

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# Slide 2: 総合評価
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, COLOR_WHITE)
slide_header(slide, "1. 総合評価", "各評価項目のサマリー")

rows = [
    ["評価項目", "評価", "コメント"],
    ["事業構想・ビジョン", "A", "明確で市場ニーズに合致"],
    ["市場分析", "B+", "概要は良いが、競合分析の具体性に課題"],
    ["プロダクト戦略", "A-", "フェーズ分けが適切、技術選定も妥当"],
    ["財務計画の整合性", "B-", "複数の数値不整合あり"],
    ["ユニットエコノミクス", "C+", "CACの算出に重大な矛盾"],
    ["営業・マーケ戦略", "B", "方向性は正しいが具体性不足"],
    ["リスク管理", "B", "主要リスクはカバー、対策の具体性に改善余地"],
]
make_table(slide, Inches(0.6), Inches(1.5), Inches(12), Inches(3.5), rows,
           col_widths=[Inches(3), Inches(1.2), Inches(7.8)])

# Summary cards
add_card(slide, Inches(0.6), Inches(5.5), Inches(3.6), Inches(1.3),
         "最大の強み", "日本市場最適化", "請求書・見積書・稟議対応", COLOR_GREEN)
add_card(slide, Inches(4.5), Inches(5.5), Inches(3.6), Inches(1.3),
         "最大の課題", "CAC過小評価", "実コストと最大2.4倍の乖離", COLOR_RED)
add_card(slide, Inches(8.4), Inches(5.5), Inches(4.3), Inches(1.3),
         "総合判断", "修正3点で投資家提示可", "ARR/実収益区別、CAC再計算、調達計画整合", COLOR_PRIMARY)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# Slide 3: ARR整合性
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, COLOR_WHITE)
slide_header(slide, "2. 財務整合性チェック① — ARR vs 顧客数 x ARPU", "期末ARRと顧客構造の整合を検証")

rows = [
    ["年度", "顧客数", "ARPU（月額）", "算出ARR", "記載ARR", "判定"],
    ["1年目", "200社", "¥50,000", "1.2億円", "1.2億円", "OK ✓"],
    ["2年目", "700社", "¥57,000", "4.788億円", "4.8億円", "OK ✓"],
    ["3年目", "1,500社", "¥67,000", "12.06億円", "12億円", "OK ✓"],
]
make_table(slide, Inches(0.6), Inches(1.5), Inches(12), Inches(2.0), rows)

# Issue box
add_shape_rect(slide, Inches(0.6), Inches(3.9), Inches(12), Inches(0.06), COLOR_RED)
add_shape_rect(slide, Inches(0.6), Inches(3.96), Inches(12), Inches(3.0), RGBColor(0xFF, 0xF0, 0xF0))

add_multi_text(slide, Inches(1.0), Inches(4.1), Inches(11.2), Inches(2.8), [
    ("問題: ARRと年間実収益の混同", 18, COLOR_RED, True, PP_ALIGN.LEFT),
    ("", 6, COLOR_RED, False, PP_ALIGN.LEFT),
    ("売上合計（1年目: 1.2億円）が期末ARR（ランレート）と同額になっている。", 13, COLOR_DARK, False, PP_ALIGN.LEFT),
    ("顧客は年間を通じて段階的に獲得するため、年間の実収益はARRより大幅に低くなる。", 13, COLOR_DARK, False, PP_ALIGN.LEFT),
    ("", 6, COLOR_DARK, False, PP_ALIGN.LEFT),
    ("  月次15%成長で1年目末200社到達の場合:", 13, COLOR_DARK_GRAY, True, PP_ALIGN.LEFT),
    ("  ・期初顧客数: 約37社 → 実サブスクリプション収益: 約0.54億円（記載1.0億円の約54%）", 13, COLOR_DARK_GRAY, False, PP_ALIGN.LEFT),
    ("", 6, COLOR_DARK, False, PP_ALIGN.LEFT),
    ("提言: 売上計画にはARR（期末ランレート）と実収益（年間累積）を別行で明記すべき", 13, COLOR_PRIMARY, True, PP_ALIGN.LEFT),
])

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# Slide 4: 損益・キャッシュフロー整合
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, COLOR_WHITE)
slide_header(slide, "2. 財務整合性チェック② — 損益計画とキャッシュフロー", "粗利益・営業利益・キャッシュフローの計算検証")

# P&L Table
add_textbox(slide, Inches(0.6), Inches(1.4), Inches(5), Inches(0.4),
            "損益計画", font_size=14, color=COLOR_DARK, bold=True)
rows_pl = [
    ["項目", "1年目", "2年目", "3年目", "判定"],
    ["売上", "1.2億", "4.8億", "12.0億", "—"],
    ["売上原価", "0.3億", "0.6億", "1.2億", "—"],
    ["粗利益", "0.9億", "4.2億", "10.8億", "OK ✓"],
    ["粗利率", "75%", "87.5%", "90%", "OK ✓"],
    ["営業費用", "3.3億", "6.5億", "10.8億", "OK ✓"],
    ["営業利益", "-2.4億", "-2.3億", "0億", "OK ✓"],
]
make_table(slide, Inches(0.6), Inches(1.8), Inches(6), Inches(3.2), rows_pl)

# CF Table
add_textbox(slide, Inches(7), Inches(1.4), Inches(5), Inches(0.4),
            "キャッシュフロー", font_size=14, color=COLOR_DARK, bold=True)
rows_cf = [
    ["項目", "1年目", "2年目", "3年目"],
    ["期首キャッシュ", "1.0億", "3.6億", "6.3億"],
    ["営業CF", "-2.4億", "-2.3億", "0億"],
    ["資金調達", "5.0億", "5.0億", "15.0億"],
    ["期末キャッシュ", "3.6億", "6.3億", "21.3億"],
]
make_table(slide, Inches(7), Inches(1.8), Inches(5.7), Inches(2.4), rows_cf)

# Issue: 2nd year funding
add_shape_rect(slide, Inches(7), Inches(4.6), Inches(5.7), Inches(0.04), COLOR_ORANGE)
add_shape_rect(slide, Inches(7), Inches(4.64), Inches(5.7), Inches(1.5), RGBColor(0xFF, 0xF8, 0xEC))
add_multi_text(slide, Inches(7.3), Inches(4.7), Inches(5.1), Inches(1.3), [
    ("注意: 2年目の調達5.0億円の出所不明", 13, COLOR_ORANGE, True, PP_ALIGN.LEFT),
    ("調達計画にはシリーズA（1年目末5億）と", 11, COLOR_DARK, False, PP_ALIGN.LEFT),
    ("シリーズB（3年目15億）のみ記載。", 11, COLOR_DARK, False, PP_ALIGN.LEFT),
    ("2年目の5億円が未説明。", 11, COLOR_DARK, False, PP_ALIGN.LEFT),
])

# Green check box
add_shape_rect(slide, Inches(0.6), Inches(5.4), Inches(6), Inches(1.5), RGBColor(0xEC, 0xFA, 0xF1))
add_multi_text(slide, Inches(0.9), Inches(5.5), Inches(5.4), Inches(1.2), [
    ("計算整合: すべてOK", 14, COLOR_GREEN, True, PP_ALIGN.LEFT),
    ("・費用内訳の合計 → 費用合計: 整合", 11, COLOR_DARK, False, PP_ALIGN.LEFT),
    ("・粗利益 = 売上 - 売上原価: 整合", 11, COLOR_DARK, False, PP_ALIGN.LEFT),
    ("・営業利益 = 粗利益 - 営業費用: 整合", 11, COLOR_DARK, False, PP_ALIGN.LEFT),
    ("・CF = 期首 + 営業CF + 調達: 整合", 11, COLOR_DARK, False, PP_ALIGN.LEFT),
    ("・人件費単価 1,000万円/人（全年度）: 妥当", 11, COLOR_DARK, False, PP_ALIGN.LEFT),
])

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# Slide 5: CAC分析（重大な不整合）
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, COLOR_WHITE)
slide_header(slide, "3. ユニットエコノミクス — CACの重大な不整合", "本レポートで最も重要な指摘事項")

rows_cac = [
    ["項目", "1年目", "2年目", "3年目"],
    ["マーケティング費用", "0.8億円", "1.5億円", "2.4億円"],
    ["セールス人件費（推定）", "0.3億円", "0.8億円", "1.5億円"],
    ["獲得費用合計", "1.1億円", "2.3億円", "3.9億円"],
    ["新規顧客数", "200社", "500社", "800社"],
    ["実質CAC", "55万円", "46万円", "49万円"],
    ["記載CAC", "30万円", "25万円", "20万円"],
    ["乖離率", "+83%", "+84%", "+144%"],
]
make_table(slide, Inches(0.6), Inches(1.5), Inches(6.5), Inches(3.8), rows_cac)

# Impact cards
add_textbox(slide, Inches(7.5), Inches(1.4), Inches(5), Inches(0.4),
            "修正後のユニットエコノミクス", font_size=14, color=COLOR_DARK, bold=True)

add_card(slide, Inches(7.5), Inches(1.9), Inches(2.4), Inches(1.3),
         "1年目 LTV/CAC", "2.2x", "記載: 4.0x → 大幅低下", COLOR_RED)
add_card(slide, Inches(10.1), Inches(1.9), Inches(2.4), Inches(1.3),
         "1年目 CAC回収", "11ヶ月", "記載: 6ヶ月 → ほぼ倍増", COLOR_RED)
add_card(slide, Inches(7.5), Inches(3.4), Inches(2.4), Inches(1.3),
         "2年目 LTV/CAC", "3.7x", "記載: 6.8x", COLOR_ORANGE)
add_card(slide, Inches(10.1), Inches(3.4), Inches(2.4), Inches(1.3),
         "3年目 LTV/CAC", "5.5x", "記載: 13.4x", COLOR_ORANGE)

# Bottom analysis
add_shape_rect(slide, Inches(0.6), Inches(5.7), Inches(12), Inches(0.04), COLOR_RED)
add_shape_rect(slide, Inches(0.6), Inches(5.74), Inches(12), Inches(1.2), RGBColor(0xFF, 0xF0, 0xF0))
add_multi_text(slide, Inches(1.0), Inches(5.8), Inches(11.2), Inches(1.0), [
    ("影響と提言", 14, COLOR_RED, True, PP_ALIGN.LEFT),
    ("修正後も2年目以降はLTV/CAC > 3xを達成可能だが、1年目は2.2xと投資効率に課題が残る。", 12, COLOR_DARK, False, PP_ALIGN.LEFT),
    ("CACは実際のマーケティング費用＋セールス人件費から算出し直すか、PLG経由の自然獲得分を除外した「有償CAC」として明記すべき。", 12, COLOR_DARK, False, PP_ALIGN.LEFT),
])

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# Slide 6: LTV検証
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, COLOR_WHITE)
slide_header(slide, "3. ユニットエコノミクス — LTV・解約率・CAC回収期間", "LTV算出の妥当性と関連指標の検証")

rows_ltv = [
    ["年度", "月次解約率", "平均寿命", "単純LTV", "記載LTV", "推定割引率"],
    ["1年目", "3.0%", "33.3ヶ月", "167万円", "120万円", "年率約14%"],
    ["2年目", "2.0%", "50.0ヶ月", "285万円", "171万円", "年率約16%"],
    ["3年目", "1.5%", "66.7ヶ月", "447万円", "268万円", "年率約16%"],
]
make_table(slide, Inches(0.6), Inches(1.5), Inches(12), Inches(2.2), rows_ltv)

# CAC payback check
rows_payback = [
    ["指標", "1年目", "2年目", "3年目", "判定"],
    ["CAC回収期間 = CAC ÷ ARPU", "6.0ヶ月", "4.4ヶ月", "3.0ヶ月", "計算OK ✓"],
    ["LTV/CAC比率", "4.0x", "6.8x", "13.4x", "計算OK ✓"],
]
make_table(slide, Inches(0.6), Inches(4.1), Inches(12), Inches(1.2), rows_payback)

add_shape_rect(slide, Inches(0.6), Inches(5.7), Inches(12), Inches(1.3), RGBColor(0xEC, 0xFA, 0xF1))
add_multi_text(slide, Inches(1.0), Inches(5.8), Inches(11.2), Inches(1.1), [
    ("評価", 14, COLOR_GREEN, True, PP_ALIGN.LEFT),
    ("LTVはディスカウントレート込みで保守的に計算されており、算出方法自体は妥当。", 12, COLOR_DARK, False, PP_ALIGN.LEFT),
    ("ただし、解約率を3.0%→1.5%へ改善する具体的なCS施策のKPI（ヘルススコア閾値、介入タイミング等）が未記載。", 12, COLOR_DARK, False, PP_ALIGN.LEFT),
    ("NRR 105%→125%を支えるアップセル/クロスセル施策とプラン移行率の想定も追加が必要。", 12, COLOR_DARK, False, PP_ALIGN.LEFT),
])

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# Slide 7: 営業・マーケティング戦略
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, COLOR_WHITE)
slide_header(slide, "4. 営業・マーケティング戦略分析", "GTM戦略、予算配分、セールス生産性")

# Strengths
add_shape_rect(slide, Inches(0.6), Inches(1.5), Inches(5.8), Inches(2.3), RGBColor(0xEC, 0xFA, 0xF1))
add_multi_text(slide, Inches(0.9), Inches(1.6), Inches(5.2), Inches(2.0), [
    ("GTM戦略の強み", 14, COLOR_GREEN, True, PP_ALIGN.LEFT),
    ("・PLG→セールス→エンタープライズの段階的展開はSaaSの王道", 11, COLOR_DARK, False, PP_ALIGN.LEFT),
    ("・パートナーチャネル（SIer、税理士事務所）は中小企業リーチに有効", 11, COLOR_DARK, False, PP_ALIGN.LEFT),
    ("・フリープランによるリード獲得は正しいアプローチ", 11, COLOR_DARK, False, PP_ALIGN.LEFT),
    ("・デジタル広告の段階的縮小（40%→20%）はPLG成熟に伴う自然な移行", 11, COLOR_DARK, False, PP_ALIGN.LEFT),
])

# Weaknesses
add_shape_rect(slide, Inches(6.8), Inches(1.5), Inches(5.9), Inches(2.3), RGBColor(0xFF, 0xF0, 0xF0))
add_multi_text(slide, Inches(7.1), Inches(1.6), Inches(5.3), Inches(2.0), [
    ("課題・不足点", 14, COLOR_RED, True, PP_ALIGN.LEFT),
    ("・Free→Paid転換率の目標が未記載（一般的には2-5%）", 11, COLOR_DARK, False, PP_ALIGN.LEFT),
    ("・フリープランの必要ユーザー獲得数が不明", 11, COLOR_DARK, False, PP_ALIGN.LEFT),
    ("・PLGを謳いつつ1年目デジタル広告40%はセールス寄り", 11, COLOR_DARK, False, PP_ALIGN.LEFT),
    ("・パートナーの立ち上げ時期と具体的提携数の目標なし", 11, COLOR_DARK, False, PP_ALIGN.LEFT),
])

# Sales productivity
add_textbox(slide, Inches(0.6), Inches(4.1), Inches(6), Inches(0.4),
            "セールスチーム生産性", font_size=14, color=COLOR_DARK, bold=True)
rows_sales = [
    ["指標", "1年目", "2年目", "3年目"],
    ["セールス人数", "3名", "8名", "15名"],
    ["新規顧客数", "200社", "500社", "800社"],
    ["1人あたり獲得数/年", "67社", "63社", "53社"],
    ["1人あたり獲得数/月", "5.6社", "5.2社", "4.4社"],
]
make_table(slide, Inches(0.6), Inches(4.5), Inches(6), Inches(2.2), rows_sales)

add_shape_rect(slide, Inches(7), Inches(4.5), Inches(5.7), Inches(2.2), RGBColor(0xF0, 0xF2, 0xF5))
add_multi_text(slide, Inches(7.3), Inches(4.6), Inches(5.1), Inches(2.0), [
    ("分析コメント", 14, COLOR_DARK, True, PP_ALIGN.LEFT),
    ("", 4, COLOR_DARK, False, PP_ALIGN.LEFT),
    ("1年目の月5.6社/人はSMB向けSaaSとして", 12, COLOR_DARK, False, PP_ALIGN.LEFT),
    ("高い水準。PLGインバウンドリードの", 12, COLOR_DARK, False, PP_ALIGN.LEFT),
    ("十分な供給が前提条件。", 12, COLOR_DARK, False, PP_ALIGN.LEFT),
    ("", 4, COLOR_DARK, False, PP_ALIGN.LEFT),
    ("3年目の生産性低下はエンタープライズ案件", 12, COLOR_DARK, False, PP_ALIGN.LEFT),
    ("増加によるものと推測され、現実的。", 12, COLOR_DARK, False, PP_ALIGN.LEFT),
])

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# Slide 8: 競合・市場ポジショニング
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, COLOR_WHITE)
slide_header(slide, "5. 競合・市場ポジショニング評価", "市場機会の妥当性と差別化の持続性")

# Market size cards
add_card(slide, Inches(0.6), Inches(1.4), Inches(2.8), Inches(1.3),
         "TAM（国内SaaS市場）", "1.8兆円", "年率13%成長", COLOR_PRIMARY)
add_card(slide, Inches(3.7), Inches(1.4), Inches(2.8), Inches(1.3),
         "SAM（10-300名企業）", "800億円", "ターゲット定義明確", COLOR_PRIMARY)
add_card(slide, Inches(6.8), Inches(1.4), Inches(2.8), Inches(1.3),
         "3年目ARR/SAM比", "1.5%", "達成可能な範囲", COLOR_GREEN)
add_card(slide, Inches(9.9), Inches(1.4), Inches(2.8), Inches(1.3),
         "年間市場成長率", "13%", "追い風の市場環境", COLOR_GREEN)

# Differentiation durability
add_textbox(slide, Inches(0.6), Inches(3.1), Inches(6), Inches(0.4),
            "差別化要素の持続性評価", font_size=14, color=COLOR_DARK, bold=True)
rows_diff = [
    ["差別化要素", "持続性", "コメント"],
    ["オールインワン統合", "中", "大手が統合を進めると優位性低下"],
    ["AI搭載", "低〜中", "全社がAI統合を進めており差別化は短命"],
    ["価格競争力", "中", "価格競争は消耗戦リスク"],
    ["日本市場最適化", "高", "商習慣対応は参入障壁として有効"],
]
make_table(slide, Inches(0.6), Inches(3.5), Inches(12), Inches(2.2), rows_diff,
           col_widths=[Inches(3), Inches(1.5), Inches(7.5)])

# Competitive analysis issues
add_shape_rect(slide, Inches(0.6), Inches(6.0), Inches(12), Inches(1.1), RGBColor(0xFF, 0xF8, 0xEC))
add_multi_text(slide, Inches(1.0), Inches(6.1), Inches(11.2), Inches(0.9), [
    ("競合分析の改善点", 13, COLOR_ORANGE, True, PP_ALIGN.LEFT),
    ("・A社/B社/C社が匿名 → 投資家への説得力が弱い　・具体的な料金比較がない　・スイッチングコスト分析がない　・海外SaaS（Notion, Monday.com, HubSpot等）の日本参入リスクへの言及なし", 11, COLOR_DARK, False, PP_ALIGN.LEFT),
])

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# Slide 9: 重要指摘事項サマリー
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, COLOR_WHITE)
slide_header(slide, "6. 重要な指摘事項と改善提案", "投資家提示前に対応すべき項目")

# Critical
add_shape_rect(slide, Inches(0.6), Inches(1.4), Inches(0.08), Inches(2.0), COLOR_RED)
add_shape_rect(slide, Inches(0.68), Inches(1.4), Inches(12.0), Inches(2.0), RGBColor(0xFF, 0xF0, 0xF0))
add_multi_text(slide, Inches(1.0), Inches(1.45), Inches(11.4), Inches(1.9), [
    ("要修正 — Critical（3件）", 16, COLOR_RED, True, PP_ALIGN.LEFT),
    ("", 4, COLOR_RED, False, PP_ALIGN.LEFT),
    ("1. ARRと年間実収益の混同  — 売上計画1年目1.2億円がARR（期末ランレート）と一致するが、実収益は約0.5-0.7億円。ARRと実収益を別行で表記", 12, COLOR_DARK, False, PP_ALIGN.LEFT),
    ("2. CACの過小評価  — 記載CACとマーケ・セールス費用の逆算値に最大2.4倍の乖離。実コスト構造に基づくCAC再計算が必要", 12, COLOR_DARK, False, PP_ALIGN.LEFT),
    ("3. 2年目キャッシュフローの調達額  — CF表に2年目5.0億円の調達記載があるが、調達計画には該当ラウンドなし。出所を明記すべき", 12, COLOR_DARK, False, PP_ALIGN.LEFT),
])

# Important
add_shape_rect(slide, Inches(0.6), Inches(3.7), Inches(0.08), Inches(2.4), COLOR_ORANGE)
add_shape_rect(slide, Inches(0.68), Inches(3.7), Inches(12.0), Inches(2.4), RGBColor(0xFF, 0xF8, 0xEC))
add_multi_text(slide, Inches(1.0), Inches(3.75), Inches(11.4), Inches(2.3), [
    ("要改善 — Important（5件）", 16, COLOR_ORANGE, True, PP_ALIGN.LEFT),
    ("", 4, COLOR_ORANGE, False, PP_ALIGN.LEFT),
    ("4. フリープラン転換率の未設定  — PLG戦略の基盤指標として、Free→Paid転換率（例: 3-5%）と必要フリーユーザー数を明記", 12, COLOR_DARK, False, PP_ALIGN.LEFT),
    ("5. NRRの達成根拠不足  — NRR 105%→125%を支えるアップセル/クロスセル施策とプラン移行率の想定を追加", 12, COLOR_DARK, False, PP_ALIGN.LEFT),
    ("6. 解約率改善の根拠不足  — 月次3.0%→1.5%への改善を支えるCS施策のKPIを追加", 12, COLOR_DARK, False, PP_ALIGN.LEFT),
    ("7. 競合の実名記載  — 投資家向け資料では具体的な競合名と定量的比較が必要", 12, COLOR_DARK, False, PP_ALIGN.LEFT),
    ("8. シリーズB 15億円の使途不明  — 損益分岐点到達後の調達意図（海外展開、M&A等）を明記", 12, COLOR_DARK, False, PP_ALIGN.LEFT),
])

# Recommended additions
add_shape_rect(slide, Inches(0.6), Inches(6.4), Inches(0.08), Inches(0.8), COLOR_PRIMARY)
add_shape_rect(slide, Inches(0.68), Inches(6.4), Inches(12.0), Inches(0.8), RGBColor(0xEE, 0xF2, 0xFF))
add_multi_text(slide, Inches(1.0), Inches(6.45), Inches(11.4), Inches(0.7), [
    ("推奨追加項目", 14, COLOR_PRIMARY, True, PP_ALIGN.LEFT),
    ("コホート分析フレームワーク ／ セールスファネル転換率 ／ NPS・CSAT目標 ／ 海外SaaS競合動向 ／ 技術的参入障壁", 12, COLOR_DARK, False, PP_ALIGN.LEFT),
])

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# Slide 10: 総括
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, COLOR_DARK)
add_shape_rect(slide, Inches(0), Inches(3.5), Inches(13.333), Inches(0.06), COLOR_ACCENT)

add_textbox(slide, Inches(1), Inches(0.8), Inches(11), Inches(0.7),
            "7. 総括", font_size=32, color=COLOR_WHITE, bold=True)

add_multi_text(slide, Inches(1), Inches(1.7), Inches(11), Inches(1.5), [
    ("中小企業向け統合SaaSという市場ニーズに即した魅力的な事業構想を持ち、", 16, COLOR_WHITE, False, PP_ALIGN.LEFT),
    ("プロダクトのフェーズ分けや技術選定は適切です。", 16, COLOR_WHITE, False, PP_ALIGN.LEFT),
    ("投資家・ステークホルダーへの提示前に、以下の3点の修正が必須です。", 16, COLOR_ACCENT, True, PP_ALIGN.LEFT),
])

# 3 action items
for i, (num, title, desc) in enumerate([
    ("1", "ARRと実収益の区別を明確化", "売上計画にARR（期末ランレート）と\n年間実収益を別行で記載し、信頼性を確保"),
    ("2", "CACを実コスト構造に基づいて再計算", "マーケ費用＋セールス人件費から算出し\nユニットエコノミクス全体を修正"),
    ("3", "2年目の資金調達計画の矛盾を解消", "CF表の2年目調達5億円の出所を明記、\nまたは調達計画にラウンドを追加"),
]):
    left = Inches(1 + i * 3.8)
    add_shape_rect(slide, left, Inches(4.0), Inches(3.4), Inches(2.5), RGBColor(0x25, 0x25, 0x3F))
    add_textbox(slide, left + Inches(0.15), Inches(4.1), Inches(0.5), Inches(0.5),
                num, font_size=28, color=COLOR_ACCENT, bold=True)
    add_textbox(slide, left + Inches(0.5), Inches(4.15), Inches(2.7), Inches(0.4),
                title, font_size=13, color=COLOR_WHITE, bold=True)
    add_textbox(slide, left + Inches(0.2), Inches(4.7), Inches(3.0), Inches(1.5),
                desc, font_size=11, color=COLOR_GRAY, bold=False)

add_textbox(slide, Inches(1), Inches(6.8), Inches(11), Inches(0.4),
            "これらを修正した上で、PLGの転換率想定やNRR達成施策を補強することで、投資家に対する説得力のある計画書になります。",
            font_size=13, color=COLOR_WHITE, bold=False)


# ── Save ──
output_path = "/home/user/claude/sales-material-analysis.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
